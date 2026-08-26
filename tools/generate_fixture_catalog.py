"""Generate deterministic, project-owned S0 fixture files and their catalog."""

from __future__ import annotations

import hashlib
import io
import json
import mimetypes
import os
import re
import struct
import zipfile
import zlib
from collections.abc import Callable
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

from docx import Document
from openpyxl import Workbook
from PIL import Image, ImageDraw
from pptx import Presentation
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen.canvas import Canvas

ROOT = Path(__file__).resolve().parents[1]
FIXTURE_ROOT = ROOT / "fixtures" / "v1"
DOCUMENT_ROOT = FIXTURE_ROOT / "documents"
RAW_ROOT = FIXTURE_ROOT / "raw-inspection"
TEMPLATE_ROOT = FIXTURE_ROOT / "templates"
FIXED_TIME = datetime(2026, 8, 21, 12, 0, 0, tzinfo=UTC)
FIXED_ZIP_TIME = (2026, 8, 21, 12, 0, 0)
CANONICAL_ZIP_CREATE_SYSTEM = 3
CANONICAL_ZIP_EXTERNAL_ATTR = 0o600 << 16
PNG_SIGNATURE = b"\x89PNG\r\n\x1a\n"
TENANT_ID = "00000000-0000-4000-8000-000000000001"
PROJECT_ID = "00000000-0000-4000-8000-000000000002"

METHODS: dict[str, tuple[str, str, float]] = {
    "UT": ("ultrasonic", "waveform_amplitude", 52.0),
    "GPR": ("ground_penetrating_radar", "two_way_time_ns", 8.0),
    "IE": ("impact_echo", "peak_frequency_hz", 4100.0),
    "RT": ("rebound_testing", "rebound_index", 36.0),
    "AE": ("acoustic_emission", "event_count", 120.0),
    "MV": ("machine_vision", "crack_width_mm", 0.2),
}


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def write_json(path: Path, value: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(
        json.dumps(value, indent=2, sort_keys=True) + "\n",
        encoding="utf-8",
        newline="\n",
    )


def write_bytes_if_changed(path: Path, content: bytes) -> None:
    """Publish deterministic bytes through one same-directory working file."""

    if path.is_file() and path.read_bytes() == content:
        return
    working_path = path.with_name(f".{path.name}.tmp")
    try:
        working_path.write_bytes(content)
        os.replace(working_path, path)
    finally:
        working_path.unlink(missing_ok=True)


def normalize_zip(content: bytes) -> bytes:
    """Return an Office ZIP container with canonical metadata and ordering."""

    with zipfile.ZipFile(io.BytesIO(content), "r") as source:
        members = [(item, source.read(item.filename)) for item in source.infolist()]
    output = io.BytesIO()
    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9) as target:
        for original, data in sorted(members, key=lambda member: member[0].filename):
            if original.filename == "docProps/core.xml":
                data = re.sub(
                    rb"(<dcterms:modified[^>]*>)[^<]*(</dcterms:modified>)",
                    rb"\g<1>2026-08-21T12:00:00Z\g<2>",
                    data,
                )
            item = zipfile.ZipInfo(original.filename, FIXED_ZIP_TIME)
            item.compress_type = zipfile.ZIP_DEFLATED
            item.create_system = CANONICAL_ZIP_CREATE_SYSTEM
            item.external_attr = CANONICAL_ZIP_EXTERNAL_ATTR
            target.writestr(item, data)
    return output.getvalue()


def png_chunk(chunk_type: bytes, content: bytes) -> bytes:
    checksum = zlib.crc32(chunk_type + content) & 0xFFFFFFFF
    return struct.pack(">I", len(content)) + chunk_type + content + struct.pack(">I", checksum)


def stored_zlib_stream(content: bytes) -> bytes:
    """Return a zlib stream using only deterministic uncompressed DEFLATE blocks."""

    output = bytearray(b"\x78\x01")
    for offset in range(0, len(content), 65535):
        block = content[offset : offset + 65535]
        final = offset + len(block) == len(content)
        output.append(1 if final else 0)
        output.extend(struct.pack("<HH", len(block), len(block) ^ 0xFFFF))
        output.extend(block)
    output.extend(struct.pack(">I", zlib.adler32(content) & 0xFFFFFFFF))
    return bytes(output)


def deterministic_png(image: Image.Image) -> bytes:
    monochrome = image.convert("1", dither=Image.Dither.NONE)
    width, height = monochrome.size
    row_size = (width + 7) // 8
    pixels = monochrome.tobytes()
    if len(pixels) != row_size * height:
        raise ValueError("unexpected packed monochrome image size")
    scanlines = b"".join(
        b"\x00" + pixels[offset : offset + row_size] for offset in range(0, len(pixels), row_size)
    )
    header = struct.pack(">IIBBBBB", width, height, 1, 0, 0, 0, 0)
    return b"".join(
        (
            PNG_SIGNATURE,
            png_chunk(b"IHDR", header),
            png_chunk(b"IDAT", stored_zlib_stream(scanlines)),
            png_chunk(b"IEND", b""),
        )
    )


def document_name(kind: str, index: int, suffix: str) -> str:
    prefix = "case" if index % 6 else "case-cn"
    return f"{prefix}-{kind}-{index:03d}{suffix}"


def create_pdf_text(path: Path, index: int) -> None:
    canvas = Canvas(str(path), invariant=1, pageCompression=1)
    canvas.setTitle(f"Synthetic NDT text document {index}")
    canvas.drawString(72, 760, f"Synthetic NDT parser fixture {index}")
    canvas.drawString(72, 735, "Structure: bridge | Method: ultrasonic | Rights: project synthetic")
    canvas.drawString(72, 710, "Table: Location A | Reading 12.5 | Unit mm")
    canvas.showPage()
    canvas.save()


def scan_image(index: int) -> Image.Image:
    image = Image.new("RGB", (900, 1200), "white")
    draw = ImageDraw.Draw(image)
    draw.text((70, 100), f"SYNTHETIC SCAN {index:03d}", fill="black")
    draw.text((70, 160), "NDT observation table", fill="black")
    draw.text((70, 220), "A-01 | 12.5 mm | review required", fill="black")
    return image


def create_pdf_scan(path: Path, index: int) -> None:
    image = scan_image(index)
    canvas = Canvas(str(path), invariant=1, pageCompression=1)
    canvas.setTitle(f"Synthetic NDT scanned PDF {index}")
    canvas.drawImage(ImageReader(image), 72, 72, width=468, height=624)
    canvas.showPage()
    canvas.save()


def create_docx(path: Path, index: int) -> None:
    document = Document()
    document.core_properties.title = f"Synthetic NDT DOCX {index}"
    document.core_properties.created = FIXED_TIME.replace(tzinfo=None)
    document.core_properties.modified = FIXED_TIME.replace(tzinfo=None)
    document.add_heading(f"Synthetic inspection note {index}", level=1)
    document.add_paragraph("Project-owned synthetic data. Not a formal engineering conclusion.")
    if index % 4 == 0:
        document.add_paragraph("Chinese text fixture: bridge inspection, ultrasonic test.")
    table = document.add_table(rows=2, cols=3)
    table.cell(0, 0).text = "Location"
    table.cell(0, 1).text = "Value"
    table.cell(0, 2).text = "Unit"
    table.cell(1, 0).text = "A-01"
    table.cell(1, 1).text = "12.5"
    table.cell(1, 2).text = "mm"
    output = io.BytesIO()
    document.save(output)
    write_bytes_if_changed(path, normalize_zip(output.getvalue()))


def create_xlsx(path: Path, index: int) -> None:
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "Synthetic observations"
    sheet.append(["Location", "Value", "Unit"])
    sheet.append(["A-01", 10 + index / 10, "mm"])
    sheet.append(["A-02", 11 + index / 10, "mm"])
    sheet.append(["A-03", 12 + index / 10, "mm"])
    sheet.append(["Total", "=SUM(B2:B4)", "mm"])
    workbook.properties.title = f"Synthetic NDT XLSX {index}"
    workbook.properties.created = FIXED_TIME.replace(tzinfo=None)
    workbook.properties.modified = FIXED_TIME.replace(tzinfo=None)
    output = io.BytesIO()
    workbook.save(output)
    write_bytes_if_changed(path, normalize_zip(output.getvalue()))


def create_pptx(path: Path, index: int) -> None:
    presentation = Presentation()
    presentation.core_properties.title = f"Synthetic NDT PPTX {index}"
    presentation.core_properties.created = FIXED_TIME.replace(tzinfo=None)
    presentation.core_properties.modified = FIXED_TIME.replace(tzinfo=None)
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = f"Synthetic inspection summary {index}"
    slide.placeholders[
        1
    ].text = "Bridge / ultrasonic\nObservation A-01: 12.5 mm\nSynthetic, review required"
    output = io.BytesIO()
    presentation.save(output)
    write_bytes_if_changed(path, normalize_zip(output.getvalue()))


def create_md(path: Path, index: int) -> None:
    content = (
        f"# Synthetic inspection note {index}\n\n"
        "Project-owned synthetic data. Not a formal engineering conclusion.\n\n"
        "| Location | Value | Unit |\n|---|---:|---|\n| A-01 | 12.5 | mm |\n"
    )
    if index % 4 == 0:
        content += "\nUnicode fixture: Chinese path and content are tested separately.\n"
    path.write_text(content, encoding="utf-8", newline="\n")


def create_txt(path: Path, index: int) -> None:
    content = (
        f"Synthetic inspection note {index}\n"
        "Structure=bridge\nMethod=ultrasonic\nLocation=A-01\nValue=12.5\nUnit=mm\n"
    )
    path.write_text(content, encoding="utf-8", newline="\n")


def create_png_scan(path: Path, index: int) -> None:
    write_bytes_if_changed(path, deterministic_png(scan_image(index)))


def media_type(path: Path) -> str:
    return mimetypes.guess_type(path.name)[0] or "application/octet-stream"


def file_record(
    path: Path, fixture_id: str, fixture_type: str, features: list[str]
) -> dict[str, Any]:
    return {
        "classification": "INTERNAL",
        "deidentification": "SYNTHETIC_NO_PERSONAL_DATA",
        "features": features,
        "fixture_id": fixture_id,
        "fixture_type": fixture_type,
        "media_type": media_type(path),
        "path": path.relative_to(ROOT).as_posix(),
        "rights_basis": "PROJECT_GENERATED_SYNTHETIC",
        "sha256": sha256(path),
        "size_bytes": path.stat().st_size,
        "training_use": "PROHIBITED",
    }


def generate_documents() -> list[dict[str, Any]]:
    creators: list[tuple[str, str, Callable[[Path, int], None], list[str]]] = [
        ("pdf-text", ".pdf", create_pdf_text, ["text", "table"]),
        ("pdf-scan", ".pdf", create_pdf_scan, ["scan", "image", "table"]),
        ("docx", ".docx", create_docx, ["text", "table", "unicode"]),
        ("xlsx", ".xlsx", create_xlsx, ["table", "formula"]),
        ("pptx", ".pptx", create_pptx, ["text", "layout"]),
        ("markdown", ".md", create_md, ["text", "table", "unicode"]),
        ("text", ".txt", create_txt, ["text"]),
        ("png-scan", ".png", create_png_scan, ["scan", "image"]),
    ]
    records: list[dict[str, Any]] = []
    for kind, suffix, creator, features in creators:
        folder = DOCUMENT_ROOT / kind
        folder.mkdir(parents=True, exist_ok=True)
        for index in range(1, 25):
            path = folder / document_name(kind, index, suffix)
            creator(path, index)
            records.append(file_record(path, f"DOC-{kind.upper()}-{index:03d}", kind, features))
    return records


def generate_raw_inspection() -> list[dict[str, Any]]:
    records: list[dict[str, Any]] = []
    for method_code, (method_name, measurement_name, baseline) in METHODS.items():
        folder = RAW_ROOT / method_code.lower()
        folder.mkdir(parents=True, exist_ok=True)
        for index in range(1, 11):
            sample_id = f"RAW-{method_code}-{index:03d}"
            path = folder / f"{sample_id.lower()}.json"
            value = baseline + (index - 5) * baseline * 0.01
            write_json(
                path,
                {
                    "classification": "INTERNAL",
                    "device": {
                        "calibration_id": "SYNTHETIC-CALIBRATION-V1",
                        "device_id": "SYNTHETIC-DEVICE",
                        "mode": "SIMULATED",
                    },
                    "location": {"component_id": "SYNTHETIC-COMPONENT-001", "station": index},
                    "measurement": {
                        "name": measurement_name,
                        "unit": "method_defined",
                        "value": value,
                    },
                    "method_code": method_code,
                    "method_name": method_name,
                    "project_id": PROJECT_ID,
                    "rights_basis": "PROJECT_GENERATED_SYNTHETIC",
                    "sample_id": sample_id,
                    "schema_version": "1.0.0",
                    "tenant_id": TENANT_ID,
                    "training_use": "PROHIBITED",
                },
            )
            records.append(file_record(path, sample_id, "raw-inspection", [method_code]))
    return records


def generate_templates() -> list[dict[str, Any]]:
    templates: dict[str, dict[str, Any]] = {
        "inspection-plan": {
            "schema_version": "1.0.0",
            "template_id": "TPL-INSPECTION-PLAN-V1",
            "required_sections": [
                "objective",
                "scope",
                "structure_and_component",
                "applicable_basis",
                "method_and_layout",
                "equipment_and_calibration",
                "procedure",
                "sampling_and_coverage",
                "acceptance_criteria",
                "safety",
                "data_management",
                "quality_control",
                "schedule",
                "deliverables",
                "limitations",
                "review_and_approval",
                "missing_input_handling",
            ],
        },
        "inspection-report": {
            "schema_version": "1.0.0",
            "template_id": "TPL-INSPECTION-REPORT-V1",
            "required_sections": [
                "identity",
                "scope",
                "plan_reference",
                "source_data",
                "method_equipment_and_calibration",
                "observations",
                "calculations_and_units",
                "figures",
                "findings",
                "limitations",
                "citations",
                "conclusion_boundary",
                "revision_history",
                "review",
                "approval",
            ],
        },
    }
    records: list[dict[str, Any]] = []
    for name, value in templates.items():
        path = TEMPLATE_ROOT / f"{name}.v1.json"
        write_json(path, value)
        records.append(file_record(path, str(value["template_id"]), "template", [name]))
    return records


def main() -> None:
    documents = generate_documents()
    raw_samples = generate_raw_inspection()
    templates = generate_templates()
    catalog = {
        "catalog_version": "1.0.0",
        "created_at": FIXED_TIME.isoformat(),
        "documents": documents,
        "external_sources": [],
        "known_gaps": [
            {
                "gap_id": "GAP-STANDARDS-RIGHTS",
                "required_action": (
                    "Legal and Knowledge Owners approve a standards rights register."
                ),
                "state": "BLOCKING",
            },
            {
                "gap_id": "GAP-REAL-DEVICE-DATA",
                "required_action": (
                    "Domain and Data Owners provide de-identified authorized samples."
                ),
                "state": "BLOCKING",
            },
        ],
        "raw_inspection_samples": raw_samples,
        "rights_statement": (
            "Project-generated synthetic fixtures; not licensed standards or real inspection "
            "evidence."
        ),
        "state": "PARTIAL_SYNTHETIC_BASELINE",
        "templates": templates,
        "training_use": "PROHIBITED",
    }
    write_json(FIXTURE_ROOT / "catalog.json", catalog)


if __name__ == "__main__":
    main()
